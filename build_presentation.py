"""Generate the Structural 8-bit ALU presentation (PPTX)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).parent / "Structural_8bit_ALU.pptx"
ASSETS = Path(__file__).parent / "presentation_assets"

THEME_BG    = RGBColor(0x0F, 0x17, 0x2A)
THEME_PANEL = RGBColor(0x1B, 0x26, 0x40)
THEME_ACCENT= RGBColor(0x4A, 0x90, 0xE2)
THEME_TEXT  = RGBColor(0xEA, 0xEE, 0xF7)
THEME_MUTED = RGBColor(0xA8, 0xB2, 0xC8)
CODE_BG     = RGBColor(0x10, 0x16, 0x26)
CODE_FG     = RGBColor(0xCD, 0xE8, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_bg(slide, color=THEME_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=THEME_TEXT,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.9))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = THEME_PANEL
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.9), SW, Inches(0.05))
    accent.line.fill.background()
    accent.fill.solid(); accent.fill.fore_color.rgb = THEME_ACCENT
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.5),
             title, size=26, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.35),
                 subtitle, size=12, color=THEME_MUTED)


def add_bullets(slide, x, y, w, h, items, *, size=16, color=THEME_TEXT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = "• " + it
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def add_code(slide, x, y, w, h, code, *, size=11):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.adjustments[0] = 0.05
    panel.line.color.rgb = THEME_ACCENT
    panel.line.width = Pt(0.75)
    panel.fill.solid(); panel.fill.fore_color.rgb = CODE_BG
    pad = Inches(0.15)
    tb = slide.shapes.add_textbox(x + pad, y + pad, w - 2*pad, h - 2*pad)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = CODE_FG


def add_placeholder(slide, x, y, w, h, label):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.04
    box.line.color.rgb = THEME_ACCENT
    box.line.width = Pt(1.25)
    box.line.dash_style = 7  # dash
    box.fill.solid(); box.fill.fore_color.rgb = THEME_PANEL
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"[ Placeholder ]\n{label}\n(insert testbench screenshot / waveform here)"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = THEME_MUTED


def add_image(slide, x, y, w, h, filename, caption=None):
    """Embed PNG asset, fitting inside (x,y,w,h). Falls back to placeholder."""
    p = ASSETS / filename
    if not p.exists():
        add_placeholder(slide, x, y, w, h, filename)
        return
    # frame
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    frame.adjustments[0] = 0.03
    frame.line.color.rgb = THEME_ACCENT
    frame.line.width = Pt(0.75)
    frame.fill.solid(); frame.fill.fore_color.rgb = THEME_PANEL
    pad = Inches(0.08)
    pic = slide.shapes.add_picture(str(p), x + pad, y + pad,
                                    width=w - 2*pad, height=h - 2*pad)
    if caption:
        add_text(slide, x, y + h + Inches(0.05), w, Inches(0.3),
                 caption, size=10, color=THEME_MUTED, align=PP_ALIGN.CENTER)


def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             "Structural 8-bit ALU  •  Verilog implementation",
             size=10, color=THEME_MUTED)
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
             f"{page_num} / {total}", size=10, color=THEME_MUTED, align=PP_ALIGN.RIGHT)


# ----------------------- SLIDES -----------------------

slides_meta = []

# 1. Title
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), SW, Inches(2.7))
band.line.fill.background(); band.fill.solid(); band.fill.fore_color.rgb = THEME_PANEL
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), Inches(0.18), Inches(2.7))
acc.line.fill.background(); acc.fill.solid(); acc.fill.fore_color.rgb = THEME_ACCENT
add_text(s, Inches(0.8), Inches(2.7), Inches(12), Inches(1.0),
         "Structural 8-bit ALU", size=44, bold=True)
add_text(s, Inches(0.8), Inches(3.7), Inches(12), Inches(0.6),
         "A Verilog implementation: add / sub / Booth multiply / radix-2 divide",
         size=20, color=THEME_MUTED)
add_text(s, Inches(0.8), Inches(4.4), Inches(12), Inches(0.5),
         "Computer Architecture  •  University Project", size=14, color=THEME_MUTED)
add_text(s, Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
         "Repository: Structural-8bit-ALU", size=11, color=THEME_MUTED)

# 2. Agenda
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "Agenda", "What we will cover")
add_bullets(s, Inches(0.7), Inches(1.3), Inches(12), Inches(5.5), [
    "Project scope and external interface",
    "Top-level architecture and dataflow",
    "Control unit FSM (IDLE / LOAD / EXEC / DONE)",
    "Add/Subtract engine (carry-select adder + XOR word-gate)",
    "Booth radix-4 multiplier engine",
    "Radix-2 restoring divider engine",
    "Common structural building blocks",
    "Verification: testbench operations and example results",
    "Performance, latency and trade-offs",
    "Summary",
], size=18)

# 3. Scope
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "1. Scope and Functional Definition")
add_bullets(s, Inches(0.7), Inches(1.2), Inches(7.5), Inches(5.5), [
    "Two 8-bit operands: A_raw, B_raw",
    "2-bit opcode selects operation",
    "Unified 16-bit result bus",
    "Handshake protocol: start → done",
    "Status flags: carry_out, overflow (add/sub)",
    "Fully structural Verilog design",
], size=17)
# Opcode table image
add_text(s, Inches(8.5), Inches(1.2), Inches(4.5), Inches(0.4),
         "Opcode Map", size=16, bold=True, color=THEME_ACCENT)
add_code(s, Inches(8.5), Inches(1.7), Inches(4.5), Inches(3.0),
         "00 → ADD     result = A + B\n"
         "01 → SUB     result = A - B\n"
         "10 → MUL     result = A * B (signed)\n"
         "11 → DIV     result = {rem, quot}",
         size=13)

# 4. External Interface
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "2. External Interface (alu_top)")
add_code(s, Inches(0.5), Inches(1.2), Inches(8.0), Inches(5.5),
"""module alu_top (
    input             clk,
    input             rst,
    input             start,
    input  [1:0]      opcode,
    input  [7:0]      A_raw,
    input  [7:0]      B_raw,
    output wire [15:0] result,
    output            carry_out,
    output            overflow,
    output            done
);""", size=14)
add_bullets(s, Inches(8.7), Inches(1.2), Inches(4.4), Inches(5.5), [
    "clk: rising-edge sampled",
    "rst: async, active high",
    "start: 1-cycle pulse OK",
    "opcode: 2-bit selector",
    "result: encoding depends on opcode",
    "done: asserted in DONE state",
], size=14)

# 5. Top-level architecture (text + reference to .circ)
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "3. Top-Level Architecture", "alu_top integrates control + datapath")
add_bullets(s, Inches(0.7), Inches(1.2), Inches(6.0), Inches(5.5), [
    "Opcode decode → engine select",
    "Operand capture in register_8bit (×2)",
    "Dispatch start to selected engine",
    "Merge engine busy/done into FSM",
    "Merge engine results into one bus",
], size=16)
add_image(s, Inches(7.0), Inches(1.2), Inches(6.0), Inches(5.5),
          "architecture_diagram.png")

# 6. Architecture diagram placeholder
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "Architecture Diagram", "Top-level block diagram")
add_image(s, Inches(0.7), Inches(1.2), Inches(12.0), Inches(5.6),
          "architecture_diagram.png")

# 6b. Top-level structural instantiation (real alu_top.v)
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "3.1 Top-Level Structural Composition",
              "alu_top.v — module instantiations")
add_code(s, Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.8),
"""control_unit brain (
    .clk(clk), .rst(rst), .start(start),
    .exec_done(effective_done), .exec_busy(effective_busy),
    .load_en(load_en), .exec_start(exec_start), .done(done));

register_8bit reg_A (.clk(clk), .rst(rst), .load_en(load_en),
                    .data_in(A_raw), .data_out(A_internal));
register_8bit reg_B (.clk(clk), .rst(rst), .load_en(load_en),
                    .data_in(B_raw), .data_out(B_internal));

adder_substractor add_sub_unit (.clk(clk), .rst(rst), .start(exec_start),
    .op1(A_internal), .op2(B_internal),
    .enable(sel_addsub), .sub_mode(sel_sub),
    .busy(addsub_busy), .done(addsub_done),
    .c_out(addsub_cout), .overflow(addsub_overflow), .result(addsub_result));

booth_radix_4_multiplier mul_unit (.clk(clk), .rst(rst), .start(exec_start),
    .multiplicand(A_internal), .multiplier(B_internal),
    .enable(sel_mul), .busy(mul_busy), .done(mul_done), .product(mul_result));

radix2_div div_unit (.clk(clk), .reset(rst), .start(exec_start & sel_div),
    .dividend(A_internal), .divisor(B_internal),
    .quotient(div_quotient), .remainder(div_remainder),
    .ready(div_ready), .done(div_done));

// Result MUX built from two structural mux2to1 instances
mux2to1 #(.w(16)) mux_result_muldiv (.in0(mul_result),
    .in1({div_remainder, div_quotient}), .sel(opcode[0]), .out(result_muldiv));
mux2to1 #(.w(16)) mux_result        (.in0({8'b0, addsub_result}),
    .in1(result_muldiv), .sel(opcode[1]), .out(result));""", size=10)

# 7. Control Unit FSM
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "4. Control Unit FSM", "src/control/control_unit.v")
add_code(s, Inches(0.5), Inches(1.2), Inches(7.8), Inches(5.5),
"""localparam IDLE=2'b00, LOAD=2'b01,
           EXEC=2'b10, DONE=2'b11;
reg [1:0] state, next_state;

always @(posedge clk or posedge rst)
  if (rst) state <= IDLE;
  else     state <= next_state;

always @(*) case (state)
  IDLE: next_state = start     ? LOAD : IDLE;
  LOAD: next_state = EXEC;
  EXEC: next_state = exec_done ? DONE : EXEC;
  DONE: next_state = start     ? DONE : IDLE;
endcase

always @(*) begin
  load_en    = (state == LOAD);
  exec_start = (state == EXEC) && !exec_busy && !exec_done;
  done       = (state == DONE);
end""", size=12)
add_bullets(s, Inches(8.2), Inches(1.2), Inches(4.8), Inches(5.5), [
    "IDLE → wait for start",
    "LOAD → latch A_raw, B_raw",
    "EXEC → drive engine, wait for done",
    "DONE → assert top-level done",
    "Transitions on rising clk edge",
    "Async reset returns FSM to IDLE",
], size=15)

# 8. Add/Sub engine
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "5. Add / Subtract Engine", "adder_substractor.v (structural)")
add_code(s, Inches(0.5), Inches(1.2), Inches(8.0), Inches(5.5),
"""// op2_xor = op2 XOR {8{sub_mode}}
xor_wordgate #(.w(8)) gate (
    .in(op2), .bit_in(sub_mode), .out(op2_xor)
);

// 8-bit add: result = op1 + op2_xor + sub_mode
carry_select_adder add (
    .op1(op1), .op2(op2_xor),
    .c_in(sub_mode),
    .result(adder_result),
    .c_out(adder_c_out)
);

assign overflow_raw =
    (op1[7] == op2_xor[7]) &&
    (adder_result[7] != op1[7]);""", size=12)
add_bullets(s, Inches(8.7), Inches(1.2), Inches(4.5), Inches(5.5), [
    "Two's complement add/sub",
    "sub_mode = 0: A + B",
    "sub_mode = 1: A + ~B + 1",
    "Carry-select adder core",
    "1-cycle done pulse",
    "Status: carry, signed overflow",
], size=14)

# 9. Booth Multiplier
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "6. Booth Radix-4 Multiplier", "booth_radix_4_multiplier.v (structural add path)")
add_code(s, Inches(0.5), Inches(1.2), Inches(8.0), Inches(5.5),
"""// Booth recoding selects operand and sub flag
case (y_ext[2:0])
  3'b001,3'b010: begin booth_operand=x_shift; booth_sub=0; end
  3'b011       : begin booth_operand=x2;      booth_sub=0; end
  3'b100       : begin booth_operand=x2;      booth_sub=1; end
  3'b101,3'b110: begin booth_operand=x_shift; booth_sub=1; end
  default      : begin booth_operand=16'b0;   booth_sub=0; end
endcase

// 16-bit accumulate built from two 8-bit carry-select adders
xor_wordgate #(.w(8)) inv_lo (.in(booth_operand[7:0]),
    .bit_in(booth_sub), .out(booth_operand_xor[7:0]));
xor_wordgate #(.w(8)) inv_hi (.in(booth_operand[15:8]),
    .bit_in(booth_sub), .out(booth_operand_xor[15:8]));

carry_select_adder add_lo (.op1(acc[7:0]),  .op2(booth_operand_xor[7:0]),
    .c_in(booth_sub),  .result(add_lo_result), .c_out(add_lo_cout));
carry_select_adder add_hi (.op1(acc[15:8]), .op2(booth_operand_xor[15:8]),
    .c_in(add_lo_cout),.result(add_hi_result), .c_out(add_hi_cout));""", size=10)
add_bullets(s, Inches(8.2), Inches(1.2), Inches(4.8), Inches(5.5), [
    "Signed 8×8 → 16-bit",
    "4 recoding iterations",
    "FSM: IDLE → CALC → FINISH",
    "Reduces partial products vs. shift-add",
    "Native signed arithmetic",
], size=14)

# 10. Divider
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "7. Radix-2 Restoring Divider", "Radix2.v (structural sub & restore)")
add_code(s, Inches(0.5), Inches(1.2), Inches(8.0), Inches(5.5),
"""// shifted_A = {A[W-2:0], Q[W-1]}
wire [WIDTH-1:0] shifted_A = {A[WIDTH-2:0], Q[WIDTH-1]};

// M_inv = ~M  (for two's-complement subtract)
xor_wordgate #(.w(WIDTH)) neg_m (
    .in(M), .bit_in(1'b1), .out(M_inv));

// Trial subtract: shifted_A - M  =  shifted_A + ~M + 1
carry_select_adder sub_adder (
    .op1(shifted_A), .op2(M_inv), .c_in(1'b1),
    .result(sub_result), .c_out(sub_cout));

// Restore path: sub_result + M  (used when subtract was negative)
carry_select_adder restore_adder (
    .op1(sub_result), .op2(M), .c_in(1'b0),
    .result(restore_result), .c_out());""", size=11)
add_bullets(s, Inches(8.2), Inches(1.2), Inches(4.8), Inches(5.5), [
    "Unsigned, WIDTH = 8",
    "Registers: A (rem), M (div), Q (quot)",
    "1 quotient bit per cycle (8 cycles)",
    "Result = {remainder, quotient}",
    "No divide-by-zero exception",
], size=14)

# 11. Common blocks
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "8. Common Structural Blocks", "src/common/")
add_bullets(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.5), [
    "full_adder_cell — 1-bit full adder primitive",
    "ripple_carry_adder — parameterized ripple chain",
    "adder_level — dual-carry precompute/select stage",
    "carry_select_adder — 8-bit ripple + carry-select",
    "xor_wordgate — word-wise XOR with replicated control",
    "mux2to1 — generic 2:1 mux",
    "register_8bit — operand storage register",
], size=14)
add_code(s, Inches(6.7), Inches(1.2), Inches(6.3), Inches(5.5),
"""// carry_select_adder.v (8-bit, structural)
module carry_select_adder(
    input  [7:0] op1, op2,
    input        c_in,
    output [7:0] result,
    output       c_out);
  wire c_mid;

  ripple_carry_adder #(.w(4)) least_significand_part (
    .op1(op1[3:0]), .op2(op2[3:0]),
    .c_in(c_in),
    .result(result[3:0]), .c_out(c_mid));

  adder_level most_significand_part (
    .op1(op1[7:4]), .op2(op2[7:4]),
    .c_in(c_mid),
    .result(result[7:4]), .c_out(c_out));
endmodule""", size=11)

# 12. End-to-end timing
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "9. End-to-End Functional Timing")
add_bullets(s, Inches(0.7), Inches(1.2), Inches(12), Inches(5.5), [
    "1. start asserted with stable A_raw, B_raw, opcode",
    "2. Control enters LOAD → operands latched",
    "3. Control enters EXEC → exec_start to engine",
    "4. Engine asserts local done",
    "5. Control enters DONE → top-level done = 1",
    "6. Caller deasserts start → return to IDLE",
], size=17)

# 13. Latency table
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "10. Performance & Latency", "Clock period = 10 ns")
add_code(s, Inches(0.7), Inches(1.3), Inches(11.5), Inches(3.5),
"""Operation   Engine cycles   Top-level cycles   Time @10ns
ADD             1                4                40 ns
SUB             1                4                40 ns
MUL             6                9                90 ns
DIV             9               12               120 ns""", size=15)
add_bullets(s, Inches(0.7), Inches(5.0), Inches(12), Inches(2.0), [
    "Add/Sub fastest: combinational datapath + handshake",
    "Mul: 4 Booth iterations (radix-4 reduces count)",
    "Div: 1 quotient bit per cycle + restore overhead",
], size=14)

# 14. Verification - testbench overview
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "11. Verification — Integrated Testbench", "sim/alu_tb.v")
add_code(s, Inches(0.5), Inches(1.2), Inches(8.5), Inches(5.5),
"""task run_op;
  input [7:0]      in_a, in_b;
  input [1:0]      in_op;
  input [15:0]     expected;
  input [8*20-1:0] label;
  begin
    @(negedge clk);
    A_raw=in_a; B_raw=in_b; opcode=in_op; start=1'b1;
    @(negedge clk); start=1'b0;
    wait (done == 1'b1); #1;
    if (result !== expected)
      $fatal(1, "FAIL [%s]", label);
    else
      $display("PASS [%0s]: A=%0d B=%0d result=%0d",
               label, in_a, in_b, result);
    wait (done == 1'b0);
  end
endtask""", size=11)
add_bullets(s, Inches(9.2), Inches(1.2), Inches(3.8), Inches(5.5), [
    "Drives all 4 opcodes",
    "Edge cases: 0, wrap, signed",
    "Self-checking: exp vs got",
    "Halts on FAIL via $fatal",
    "Generates alu_sim.vcd",
], size=13)

# 15. Test coverage list
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "11.1 Test Vectors Overview")
add_code(s, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.5),
"""// ADD (opcode 00)
ADD 5+3          → 8
ADD 100+55       → 155
ADD 255+1 wrap   → 0
ADD 128+128 wrap → 0

// SUB (opcode 01)
SUB 10-4         → 6
SUB 0-1 wrap     → 255
SUB 200-50       → 150""", size=12)
add_code(s, Inches(6.9), Inches(1.2), Inches(6.2), Inches(5.5),
"""// MUL (opcode 10, signed Booth)
MUL 7*6          → 42
MUL 15*15        → 225
MUL 127*2        → 254
MUL -1*2 signed  → 0xFFFE

// DIV (opcode 11) → {rem, quot}
DIV 42/6         → {0, 7}
DIV 17/5         → {2, 3}
DIV 255/16       → {15, 15}
DIV 100/7        → {2, 14}""", size=12)

# 16. Example: ADD waveform placeholder
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "Example: ADD 5 + 3", "Expected result = 8")
add_code(s, Inches(0.5), Inches(1.2), Inches(5.5), Inches(2.5),
"""run_op(8'd5, 8'd3, 2'b00,
       16'd8, "ADD 5+3");

PASS [ADD 5+3]:
  A=5  B=3
  result=8 (0x0008)
  carry=0  ovf=0""", size=12)
add_image(s, Inches(6.3), Inches(1.2), Inches(6.5), Inches(5.5),
          "wave_add.png", caption="Cycle-accurate waveform of the ADD operation")

# 17. Example: SUB waveform placeholder
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "Example: SUB 10 - 4", "Expected result = 6")
add_code(s, Inches(0.5), Inches(1.2), Inches(5.5), Inches(2.5),
"""run_op(8'd10, 8'd4, 2'b01,
       16'd6, "SUB 10-4");

PASS [SUB 10-4]:
  A=10  B=4
  result=6 (0x0006)
  carry=1  ovf=0""", size=12)
add_image(s, Inches(6.3), Inches(1.2), Inches(6.5), Inches(5.5),
          "wave_sub.png", caption="Cycle-accurate waveform of the SUB operation")

# 18. Example: MUL waveform placeholder
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "Example: MUL 15 × 15", "Expected result = 225")
add_code(s, Inches(0.5), Inches(1.2), Inches(5.5), Inches(2.5),
"""run_op(8'd15, 8'd15, 2'b10,
       16'd225, "MUL 15*15");

PASS [MUL 15*15]:
  A=15  B=15
  result=225 (0x00E1)""", size=12)
add_image(s, Inches(6.3), Inches(1.2), Inches(6.5), Inches(5.5),
          "wave_mul.png", caption="Booth radix-4 multiplier (6 EXEC cycles)")

# 19. Example: DIV waveform placeholder
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "Example: DIV 100 / 7", "Expected = {rem=2, quot=14}")
add_code(s, Inches(0.5), Inches(1.2), Inches(5.5), Inches(2.5),
"""run_op(8'd100, 8'd7, 2'b11,
       {8'd2, 8'd14}, "DIV 100/7");

PASS [DIV 100/7]:
  A=100  B=7
  result=0x020E
  rem=2  quot=14""", size=12)
add_image(s, Inches(6.3), Inches(1.2), Inches(6.5), Inches(5.5),
          "wave_div.png", caption="Radix-2 restoring divider (9 EXEC cycles)")

# 20. Full simulation run placeholder
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "Simulation Run Output", "Console transcript from run.sh")
add_image(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.6),
          "console_transcript.png")

# 21. Trade-offs
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "12. Algorithm Trade-offs")
add_bullets(s, Inches(0.7), Inches(1.2), Inches(12), Inches(5.5), [
    "Add/Sub: carry-select adder — faster than ripple, more area",
    "Multiply: Booth radix-4 — fewer iterations, signed-native, multi-cycle",
    "Divide: radix-2 restoring — simple control, highest latency",
    "Overall priority: structural clarity & module reuse",
    "Trade: not minimum latency, not maximum throughput",
], size=17)

# 22. Summary
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
add_title_bar(s, "Summary")
add_bullets(s, Inches(0.7), Inches(1.2), Inches(12), Inches(5.5), [
    "Fully structural 8-bit ALU in Verilog",
    "Unified handshake (start/done) across 4 engines",
    "Clear FSM control + reusable common primitives",
    "Verified with self-checking integrated testbench",
    "Architecture captured in alu_architecture.circ (Logisim Evolution)",
], size=18)

# 23. Q&A
s = prs.slides.add_slide(BLANK); slides_meta.append(s); add_bg(s)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), SW, Inches(2.0))
band.line.fill.background(); band.fill.solid(); band.fill.fore_color.rgb = THEME_PANEL
add_text(s, 0, Inches(3.0), SW, Inches(1.0),
         "Questions?", size=54, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(4.0), SW, Inches(0.6),
         "Thank you", size=22, color=THEME_MUTED, align=PP_ALIGN.CENTER)

# Footers
total = len(slides_meta)
for idx, slide in enumerate(slides_meta, start=1):
    add_footer(slide, idx, total)

prs.save(OUT)
print(f"Saved: {OUT}")
